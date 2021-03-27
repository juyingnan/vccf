from pptx import Presentation
from skimage import io
import math
import csv


def write_csv(path, list_of_columns, list_of_names=None):
    with open(path, 'w', newline='') as csvfile:
        writer = csv.writer(csvfile, delimiter='\t', quotechar='|', quoting=csv.QUOTE_MINIMAL)
        if list_of_names is not None:
            writer.writerow(list_of_names)
        lines = []
        for _i in range(len(list_of_columns[0])):
            lines.append([column[_i] for column in list_of_columns])
        writer.writerows(lines)


# prs = Presentation(r'C:\Users\bunny\Desktop\KidneyAnnotated-GW.pptx')
# print(prs.slide_height, prs.slide_width)
# slide = prs.slides[0]
#
# shapes = slide.shapes
#
# print(len(shapes))
#
unit_per_um = 1024646 / 50
unit_per_pixel = 28346400 / 2232
pixel_per_um = unit_per_um / unit_per_pixel

nuclei_id_list = list()
nuclei_x_list = list()
nuclei_y_list = list()
nuclei_distance_list = list()
nuclei_nearest_vessel_x_list = list()
nuclei_nearest_vessel_y_list = list()
#
# for i in range(len(shapes)):
#     shape = shapes[i]
#     x = (shape.left + shape.width / 2) / unit_per_um
#     y = (shape.top + shape.height / 2) / unit_per_um
#     # print(x, y, shape.left, shape.width, shape.top, shape.height)
#     nuclei_x_list.append(x)
#     nuclei_y_list.append(y)
#     nuclei_id_list.append(i)

nuclei_image = io.imread('../images/nuclei_ml.png')  # [::-1, :]
_nid = 0
for i in range(len(nuclei_image)):
    row = nuclei_image[i]
    for j in range(len(row)):
        pixel = row[j]
        if pixel[0] > 200:
            # each nuclear is 2x2 pixels
            # get the average (+0.5)
            # and erase other 3 points
            nuclei_y_list.append((i + 0.5) / pixel_per_um)
            nuclei_x_list.append((j + 0.5) / pixel_per_um)
            nuclei_image[i][j + 1] *= 0
            nuclei_image[i + 1][j] *= 0
            nuclei_image[i + 1][j + 1] *= 0
            nuclei_id_list.append(_nid)
            _nid += 1
print(len(nuclei_x_list))

vessel_x_list = list()
vessel_y_list = list()

vessel_image = io.imread('../images/vessles.png')  # [::-1, :]
for i in range(len(vessel_image)):
    row = vessel_image[i]
    for j in range(len(row)):
        pixel = row[j]
        if pixel[0] > 200:
            vessel_y_list.append(i / pixel_per_um)
            vessel_x_list.append(j / pixel_per_um)
print(len(vessel_x_list))

for nid in nuclei_id_list:
    _min_dist = 100
    _min_vessel_x = 0
    _min_vessel_y = 0
    _nx = nuclei_x_list[nid]
    _ny = nuclei_y_list[nid]
    for v in range(len(vessel_x_list)):
        _vx = vessel_x_list[v]
        _vy = vessel_y_list[v]
        if abs(_nx - _vx) < _min_dist and abs(_ny - _vy) < _min_dist:
            _dist = math.sqrt((_nx - _vx) ** 2 + (_ny - _vy) ** 2)
            if _dist < _min_dist:
                _min_dist = _dist
                _min_vessel_x = _vx
                _min_vessel_y = _vy
    nuclei_distance_list.append(_min_dist)
    nuclei_nearest_vessel_x_list.append(_min_vessel_x)
    nuclei_nearest_vessel_y_list.append(_min_vessel_y)
    if nid % 10 == 0:
        print('\r' + str(nid), end='')

print(len(nuclei_id_list), len(nuclei_x_list), len(nuclei_y_list), len(nuclei_distance_list),
      len(nuclei_nearest_vessel_x_list), len(nuclei_nearest_vessel_y_list))

write_csv('../csv/nuclei_ml.csv',
          [nuclei_id_list,
           nuclei_x_list,
           nuclei_y_list,
           nuclei_distance_list,
           nuclei_nearest_vessel_x_list,
           nuclei_nearest_vessel_y_list],
          ['id', 'x', 'y', 'distance', 'vx', 'vy'])

write_csv('../csv/vessel.csv',
          [vessel_x_list, vessel_y_list],
          ['x', 'y'])

import matplotlib.pyplot as plt

# plt.scatter(vessel_x_list, vessel_y_list, c='r')
# plt.scatter(nuclei_x_list, nuclei_y_list, c='b')
plt.hist(nuclei_distance_list, bins=100)
plt.show()
